"""Generate a beginner-friendly AI learning module as a PPTX deck.

This deck is intentionally linear and structured with a clear learning flow (alur).
It is generated from the existing repo context (program modules + docs).

Run (Windows PowerShell):
  & "./.venv/Scripts/python.exe" make_pptx.py

Output:
  AI_Learning_Modul.pptx
"""

from __future__ import annotations

from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
OUT = ROOT / "AI_Learning_Modul.pptx"


# --- Comic theme (simple, readable, cross-platform) ---
# Font "Comic Sans MS" biasanya ada di Windows; jika tidak ada, PowerPoint akan fallback otomatis.
FONT_COMIC = "Comic Sans MS"
FONT_FALLBACK = "Calibri"

# Palette (soft background + thick black outline, comic-ish)
COLOR_BG = RGBColor(255, 253, 240)  # warm paper
COLOR_INK = RGBColor(25, 25, 25)    # near-black
COLOR_MUTED = RGBColor(80, 80, 80)

COLOR_NOTE = RGBColor(200, 230, 201)      # light green
COLOR_FLOW = RGBColor(179, 229, 252)      # light blue
COLOR_CODE_BG = RGBColor(245, 245, 245)   # light gray


def _theme_font() -> str:
    # We can't reliably detect installed fonts with python-pptx.
    # Setting Comic Sans will render as Comic Sans when available, otherwise fallback.
    return FONT_COMIC


def _apply_slide_background(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BG


def _style_text_frame(text_frame, *, bold_first_line: bool = False) -> None:
    mono_fonts = {"Consolas", "Courier New", "Menlo", "Monaco"}
    for pi, p in enumerate(text_frame.paragraphs):
        # Paragraph-level font
        p.font.name = _theme_font()
        if p.font.size is None:
            p.font.size = Pt(18)
        # Avoid reading .rgb (python-pptx raises when color type is NONE/scheme)
        p.font.color.rgb = COLOR_INK
        if bold_first_line and pi == 0:
            p.font.bold = True

        # Run-level font (overrides paragraph font in PowerPoint)
        for run in p.runs:
            if run.font.name not in mono_fonts:
                run.font.name = _theme_font()
            run.font.color.rgb = COLOR_INK


def _style_title_shape(title_shape) -> None:
    if title_shape is None:
        return
    tf = title_shape.text_frame
    _style_text_frame(tf, bold_first_line=True)
    for p in tf.paragraphs:
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = COLOR_INK


def _set_slide_title(slide, title: str) -> None:
    if slide.shapes.title is not None:
        slide.shapes.title.text = title
        _style_title_shape(slide.shapes.title)
        return

    # Fallback: add a title textbox
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12.1), Inches(0.6))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(34)
    run.font.bold = True
    run.font.name = _theme_font()
    run.font.color.rgb = COLOR_INK


def _add_bullets(slide, left: float, top: float, width: float, height: float, items: list[str]) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()

    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(20)
        p.font.name = _theme_font()
        p.font.color.rgb = COLOR_INK
        p.space_after = Pt(6)


def _add_bullets_fs(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    items: list[str],
    font_size: int = 18,
    level: int = 0,
) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()

    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = level
        p.font.size = Pt(font_size)
        p.font.name = _theme_font()
        p.font.color.rgb = COLOR_INK
        p.space_after = Pt(4)


def _add_code_box(slide, left: float, top: float, width: float, height: float, code_lines: list[str]) -> None:
    box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    box.fill.solid()
    box.fill.fore_color.rgb = COLOR_CODE_BG
    box.line.color.rgb = COLOR_INK
    box.line.width = Pt(2)

    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True

    for i, line in enumerate(code_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.name = "Consolas"
        run.font.size = Pt(15)
        run.font.color.rgb = COLOR_INK


def _add_section_slide(prs: Presentation, title: str, subtitle: str | None = None) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    _apply_slide_background(slide)
    _set_slide_title(slide, title)
    if subtitle:
        box = slide.shapes.add_textbox(Inches(0.9), Inches(2.2), Inches(12.0), Inches(2.0))
        tf = box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(22)
        p.font.name = _theme_font()
        p.font.color.rgb = COLOR_MUTED


def _add_two_columns(slide, left_title: str, left: list[str], right_title: str, right: list[str]) -> None:
    hdr = slide.shapes.add_textbox(Inches(0.9), Inches(1.7), Inches(5.8), Inches(0.4))
    tf = hdr.text_frame
    tf.text = left_title
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.name = _theme_font()
    tf.paragraphs[0].font.color.rgb = COLOR_INK

    hdr2 = slide.shapes.add_textbox(Inches(7.0), Inches(1.7), Inches(5.8), Inches(0.4))
    tf2 = hdr2.text_frame
    tf2.text = right_title
    tf2.paragraphs[0].font.size = Pt(18)
    tf2.paragraphs[0].font.bold = True
    tf2.paragraphs[0].font.name = _theme_font()
    tf2.paragraphs[0].font.color.rgb = COLOR_INK

    _add_bullets_fs(slide, 0.9, 2.2, 5.8, 4.9, left, font_size=16)
    _add_bullets_fs(slide, 7.0, 2.2, 5.8, 4.9, right, font_size=16)


def _add_note(slide, text: str) -> None:
    box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.75),
        Inches(6.6),
        Inches(12.0),
        Inches(0.75),
    )
    box.fill.solid()
    box.fill.fore_color.rgb = COLOR_NOTE
    box.line.color.rgb = COLOR_INK
    box.line.width = Pt(2)

    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(14)
    p.font.name = _theme_font()
    p.font.color.rgb = COLOR_INK


def _add_flow(slide, steps: list[str]) -> None:
    """Add a simple left-to-right flow (alur belajar)."""

    x = 0.65
    y = 2.2
    w = 2.25
    h = 0.9
    gap = 0.25

    for i, step in enumerate(steps):
        box = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(x + i * (w + gap)),
            Inches(y),
            Inches(w),
            Inches(h),
        )
        box.fill.solid()
        box.fill.fore_color.rgb = COLOR_FLOW
        box.line.color.rgb = COLOR_INK
        box.line.width = Pt(2)

        tf = box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = step
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.name = _theme_font()
        p.font.color.rgb = COLOR_INK

        # Arrow (except after last)
        if i < len(steps) - 1:
            arrow = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW,
                Inches(x + (i + 1) * (w + gap) - gap),
                Inches(y + 0.22),
                Inches(gap),
                Inches(h - 0.44),
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = COLOR_INK
            arrow.line.color.rgb = COLOR_INK


def _add_table_like(slide, title_left: str, left_items: list[str], title_right: str, right_items: list[str]) -> None:
    # Column headers
    hdr = slide.shapes.add_textbox(Inches(0.9), Inches(2.0), Inches(5.8), Inches(0.4))
    tf = hdr.text_frame
    tf.text = title_left
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.name = _theme_font()
    tf.paragraphs[0].font.color.rgb = COLOR_INK

    hdr2 = slide.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.8), Inches(0.4))
    tf2 = hdr2.text_frame
    tf2.text = title_right
    tf2.paragraphs[0].font.size = Pt(18)
    tf2.paragraphs[0].font.bold = True
    tf2.paragraphs[0].font.name = _theme_font()
    tf2.paragraphs[0].font.color.rgb = COLOR_INK

    _add_bullets(slide, 0.9, 2.5, 5.8, 3.9, left_items)
    _add_bullets(slide, 7.0, 2.5, 5.8, 3.9, right_items)


def build_deck() -> Presentation:
    prs = Presentation()
    # Set widescreen 16:9 size (PowerPoint default widescreen)
    prs.slide_width = Inches(13.3333333333)
    prs.slide_height = Inches(7.5)

    # --- Cover ---
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    _apply_slide_background(slide)
    slide.shapes.title.text = "Belajar AI untuk Pemula"
    _style_title_shape(slide.shapes.title)
    subtitle = slide.placeholders[1]
    subtitle.text = (
        "Modul latihan (Regression, Classification, NN, CNN)\n"
        "+ alur belajar + latihan + pembahasan + troubleshooting\n"
        f"Generated: {date.today().isoformat()}  |  Repo: Ai Learn"
    )
    _style_text_frame(subtitle.text_frame)

    # --- Agenda (lengkap) ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _apply_slide_background(slide)
    _set_slide_title(slide, "Agenda")
    _add_bullets_fs(
        slide,
        0.9,
        1.7,
        12.0,
        5.2,
        [
            "0) Cara pakai repo + cara run yang benar",
            "1) Alur ML: Data → Train → Evaluate → Improve",
            "2) Konsep penting: fitur/target, split, overfitting",
            "3) Metrik regression: MAE/RMSE/R² (cara baca)",
            "4) Metrik classification: accuracy/precision/recall/F1 + confusion matrix",
            "5) Modul 1–4 + latihan + pembahasan",
            "6) Troubleshooting + checklist project",
        ],
        font_size=18,
    )

    # --- File apa saja di repo ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _apply_slide_background(slide)
    _set_slide_title(slide, "Isi Repo (Yang Akan Anda Pakai)")
    _add_two_columns(
        slide,
        "Program latihan",
        [
            "01_linear_regression.py (Regression)",
            "02_classification.py (Classification)",
            "03_neural_network.py (MNIST NN)",
            "04_computer_vision.py (CNN CIFAR-10)",
        ],
        "Dokumentasi",
        [
            "BUKU_PANDUAN_AI.pdf (buku modul)",
            "AI_Learning_Modul.pptx (slide)",
            "TUTORIAL.md (cara run + tips)",
            "PANDUAN.md (perilaku normal vs error)",
            "GLOSSARY.md (istilah)",
        ],
    )
    _add_note(slide, "Prinsip: baca 1–2 halaman → jalankan 1 program → ubah 1 parameter → catat hasil.")

    # --- Flow ---
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    _set_slide_title(slide, "Alur Belajar (Flow)")
    _add_flow(slide, ["Setup", "Data", "Train", "Evaluate", "Improve", "Project"])  # type: ignore[arg-type]
    _add_bullets_fs(
        slide,
        0.9,
        3.6,
        12.0,
        2.7,
        [
            "Kalau Anda bingung: kembali ke flow ini. Semua modul mengikuti pola yang sama.",
            "Target akhir: Anda bisa menjelaskan 2 hal: (1) hasil/metric, (2) mengapa hasilnya begitu.",
        ],
        font_size=18,
    )

    # --- Setup / menjalankan ---
    _add_section_slide(prs, "0) Setup", "Pastikan cara run benar (biar tidak error palsu)")

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _set_slide_title(slide, "Cara Menjalankan (Windows PowerShell)")
    _add_bullets_fs(
        slide,
        0.9,
        1.8,
        12.0,
        2.3,
        [
            "Selalu jalankan Python dari .venv project (bukan Python global).",
            "Format perintah: & \"./.venv/Scripts/python.exe\" <nama_file.py>",
        ],
        font_size=18,
    )
    _add_code_box(
        slide,
        0.9,
        3.6,
        12.0,
        2.2,
        [
            '& "./.venv/Scripts/python.exe" 01_linear_regression.py',
            '& "./.venv/Scripts/python.exe" 02_classification.py',
            '& "./.venv/Scripts/python.exe" 03_neural_network.py',
            '& "./.venv/Scripts/python.exe" 04_computer_vision.py',
        ],
    )
    _add_note(slide, "Jika Anda pakai VS Code: pastikan interpreter VS Code juga menunjuk ke .venv.")

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _set_slide_title(slide, "Matplotlib: Grafik Muncul, Lalu Apa?")
    _add_two_columns(
        slide,
        "Perilaku normal",
        [
            "Program berhenti sejenak di plt.show()",
            "Anda melihat grafik",
            "Tutup window grafik → program lanjut",
        ],
        "Yang sering disangka error",
        [
            "Tutup window lalu muncul KeyboardInterrupt",
            "Atau Anda tekan Ctrl+C",
            "Artinya program dihentikan manual",
        ],
    )
    _add_note(slide, "Rule of thumb: kalau hasil sebelumnya sudah keluar dan Anda yang menghentikan, itu bukan bug model.")

    # --- Konsep inti ---
    _add_section_slide(prs, "1) Fondasi", "Konsep wajib sebelum masuk ke model")

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _set_slide_title(slide, "Fitur (X) dan Target (y)")
    _add_bullets_fs(
        slide,
        0.9,
        1.8,
        12.0,
        5.2,
        [
            "Fitur (X): data input yang dipakai model untuk membuat prediksi.",
            "Target (y): jawaban yang ingin dipelajari model.",
            "Regression: y berupa angka (mis. harga).",
            "Classification: y berupa kelas (mis. spam / tidak).",
            "Kesalahan umum: fitur bocor (data leakage) — memasukkan info masa depan ke fitur.",
        ],
        font_size=18,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _set_slide_title(slide, "Train / Validation / Test")
    _add_bullets_fs(
        slide,
        0.9,
        1.8,
        12.0,
        5.2,
        [
            "Train: data untuk belajar (fit).",
            "Validation: data untuk memilih setting (hyperparameter).",
            "Test: data untuk evaluasi final (simulasi dunia nyata).",
            "Kalau pemula: minimal train/test dulu, lalu tambah validation saat tuning.",
        ],
        font_size=18,
    )
    _add_note(slide, "Prinsip: test hanya dipakai sekali untuk hasil final.")

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _set_slide_title(slide, "Overfitting vs Underfitting")
    _add_two_columns(
        slide,
        "Underfitting",
        [
            "Model terlalu sederhana",
            "Train score rendah",
            "Test score rendah",
            "Solusi: model lebih kuat, fitur lebih informatif",
        ],
        "Overfitting",
        [
            "Model terlalu menghafal",
            "Train score tinggi",
            "Test score turun",
            "Solusi: regularisasi, data lebih banyak, early stopping",
        ],
    )

    # --- Metrik regression ---
    _add_section_slide(prs, "2) Metrik", "Cara membaca hasil agar tidak salah interpretasi")

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _set_slide_title(slide, "Regression Metrics: MAE / RMSE / R²")
    _add_bullets_fs(
        slide,
        0.9,
        1.8,
        12.0,
        5.2,
        [
            "MAE: rata-rata selisih |aktual - prediksi| (lebih mudah dipahami).",
            "RMSE: seperti MAE tapi menghukum error besar (lebih sensitif outlier).",
            "R²: seberapa besar variasi y yang bisa dijelaskan model (0.0–1.0).",
            "Interpretasi cepat: R² ~0.8+ bagus untuk data yang cukup linear.",
        ],
        font_size=18,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _set_slide_title(slide, "Cara Membaca R² (Contoh)")
    _add_bullets_fs(
        slide,
        0.9,
        1.8,
        12.0,
        3.2,
        [
            "R² = 0.90 → model menjelaskan 90% variasi harga (pola sangat tertangkap).",
            "R² = 0.50 → masih banyak variasi yang belum terjelaskan (fitur kurang / noise tinggi).",
            "R² < 0 → model lebih buruk dari prediksi rata-rata (indikasi masalah serius).",
        ],
        font_size=18,
    )
    _add_note(slide, "Jangan lihat 1 metric saja: cek juga plot residual/pred vs actual jika tersedia.")

    # --- Metrik classification ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _set_slide_title(slide, "Classification Metrics")
    _add_bullets_fs(
        slide,
        0.9,
        1.8,
        12.0,
        5.2,
        [
            "Accuracy: proporsi prediksi yang benar (bagus jika kelas seimbang).",
            "Precision: dari yang diprediksi positif, berapa yang benar positif.",
            "Recall: dari yang sebenarnya positif, berapa yang berhasil ditangkap.",
            "F1: rata-rata harmonik precision & recall (bagus untuk data tidak seimbang).",
        ],
        font_size=18,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _set_slide_title(slide, "Confusion Matrix (Cara Baca)")
    _add_two_columns(
        slide,
        "Makna kotak",
        [
            "TP: benar positif",
            "TN: benar negatif",
            "FP: salah positif (false alarm)",
            "FN: salah negatif (miss)",
        ],
        "Kapan fokus apa?",
        [
            "Jika false alarm mahal → perhatikan precision",
            "Jika miss mahal → perhatikan recall",
            "Jika seimbang → lihat F1",
        ],
    )
    _add_note(slide, "Iris dataset biasanya seimbang, jadi accuracy cukup representatif. Pada kasus fraud/spam, F1 lebih aman.")

    # --- Modul 1 ---
    _add_section_slide(prs, "Modul 1", "Linear Regression — prediksi angka")

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _set_slide_title(slide, "Modul 1: Tujuan & Output")
    _add_bullets_fs(
        slide,
        0.9,
        1.8,
        12.0,
        4.2,
        [
            "Tujuan: memprediksi harga rumah dari ukuran.",
            "Output: rumus garis (slope & intercept), R²/MAE/RMSE, dan plot.",
            "File: 01_linear_regression.py",
        ],
        font_size=18,
    )
    _add_code_box(slide, 0.9, 5.2, 12.0, 1.2, ['& "./.venv/Scripts/python.exe" 01_linear_regression.py'])

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _set_slide_title(slide, "Modul 1: Langkah di Kode (Flow)")
    _add_flow(slide, ["Buat data", "Plot", "Split", "Fit", "Predict", "Evaluate"])  # type: ignore[arg-type]
    _add_bullets_fs(
        slide,
        0.9,
        3.6,
        12.0,
        2.7,
        [
            "Perhatikan: bentuk data X harus 2D untuk scikit-learn.",
            "Jangan lupa: evaluasi pakai data test, bukan training.",
        ],
        font_size=18,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _set_slide_title(slide, "Modul 1: Mengapa Jawabannya Seperti Itu?")
    _add_bullets_fs(
        slide,
        0.9,
        1.8,
        12.0,
        5.2,
        [
            "Linear regression mencari garis yang meminimalkan total error kuadrat (least squares).",
            "Slope = dampak perubahan ukuran terhadap harga (per 1 m²).",
            "Intercept = harga dasar ketika ukuran mendekati 0 (makna tergantung data).",
            "Jika noise besar → R² turun karena pola linear tertutup variasi acak.",
        ],
        font_size=18,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _set_slide_title(slide, "Modul 1: Latihan")
    _add_bullets_fs(
        slide,
        0.9,
        1.8,
        12.0,
        5.2,
        [
            "Latihan A: Naikkan noise → apa yang terjadi pada RMSE dan R²?",
            "Latihan B: Tambah jumlah sampel → apakah hasil makin stabil?",
            "Latihan C: Buat outlier ekstrem → metric mana yang paling terpengaruh (MAE vs RMSE)?",
        ],
        font_size=18,
    )
    _add_note(slide, "Pembahasan singkat: noise/outlier biasanya membuat RMSE naik lebih tajam daripada MAE.")

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _set_slide_title(slide, "Modul 1: Pembahasan Latihan (Step-by-step)")
    _add_bullets_fs(
        slide,
        0.9,
        1.8,
        12.0,
        5.2,
        [
            "A) Noise naik: titik data makin menyebar → garis terbaik tetap dicari, tapi error rata-rata naik.",
            "   Dampak: RMSE naik, R² turun (variasi y lebih banyak dari noise).",
            "B) Sample size naik: estimasi slope/intercept makin stabil (lebih sedikit 'kebetulan').",
            "   Dampak: R² biasanya lebih konsisten antar run.",
            "C) Outlier ekstrem: error kuadrat membuat RMSE sangat sensitif.",
            "   Dampak: RMSE melonjak, MAE naik tapi biasanya tidak setajam RMSE.",
        ],
        font_size=17,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _set_slide_title(slide, "Modul 1: Contoh Kasus Nyata (Harga Rumah)")
    _add_two_columns(
        slide,
        "Fitur (X)",
        [
            "Luas bangunan, luas tanah",
            "Lokasi (jarak ke pusat kota)",
            "Jumlah kamar, usia bangunan",
            "Akses jalan/transportasi",
        ],
        "Target (y) + Evaluasi",
        [
            "Target: harga jual (angka)",
            "Metric: MAE/RMSE (dalam rupiah) + R²",
            "Waspada: data leakage (harga appraisal/indikator masa depan)",
            "Output berguna: interval harga & alasan fitur penting",
        ],
    )
    _add_note(slide, "Checklist: mulai dari baseline linear → tambah fitur → bandingkan MAE di test.")

    # --- Modul 2 ---
    _add_section_slide(prs, "Modul 2", "Classification — Decision Tree")

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _set_slide_title(slide, "Modul 2: Tujuan & Output")
    _add_bullets_fs(
        slide,
        0.9,
        1.8,
        12.0,
        4.2,
        [
            "Tujuan: klasifikasi bunga iris ke 3 spesies.",
            "Output: accuracy + classification report + confusion matrix + visualisasi tree.",
            "File: 02_classification.py",
        ],
        font_size=18,
    )
    _add_code_box(slide, 0.9, 5.2, 12.0, 1.2, ['& "./.venv/Scripts/python.exe" 02_classification.py'])

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _set_slide_title(slide, "Modul 2: Mengapa Decision Tree Bisa Benar?")
    _add_bullets_fs(
        slide,
        0.9,
        1.8,
        12.0,
        5.2,
        [
            "Tree membuat aturan IF-THEN yang memecah data agar tiap bagian makin 'murni'.",
            "Split dipilih yang paling meningkatkan pemisahan kelas (information gain / gini).",
            "Iris: petal length/width sering jadi fitur paling penting karena separasi jelas.",
        ],
        font_size=18,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _set_slide_title(slide, "Modul 2: Cara Baca Confusion Matrix")
    _add_bullets_fs(
        slide,
        0.9,
        1.8,
        12.0,
        5.2,
        [
            "Diagonal besar = banyak prediksi benar.",
            "Jika versicolor sering diprediksi virginica: berarti fitur overlap di area tersebut.",
            "Fokus improvement: cek fitur tambahan / model ensemble (next level).",
        ],
        font_size=18,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _set_slide_title(slide, "Modul 2: Latihan")
    _add_bullets_fs(
        slide,
        0.9,
        1.8,
        12.0,
        5.2,
        [
            "Latihan A: Ubah max_depth (1, 2, 3, 5, None). Catat train vs test accuracy.",
            "Latihan B: Cari fitur paling penting. Apakah selalu sama?",
            "Latihan C: Coba cross-validation dan bandingkan dengan 1 kali split.",
        ],
        font_size=18,
    )
    _add_note(slide, "Pembahasan singkat: depth besar → train naik, test bisa turun (overfitting).")

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _set_slide_title(slide, "Modul 2: Pembahasan Latihan (Step-by-step)")
    _add_bullets_fs(
        slide,
        0.9,
        1.8,
        12.0,
        5.2,
        [
            "A) max_depth kecil: aturan sedikit → underfitting (train & test rendah).",
            "   max_depth besar: aturan banyak → train tinggi, test bisa turun (overfitting).",
            "B) Fitur penting bisa berubah jika data split berubah sedikit (tree sensitif).",
            "   Solusi: lihat rata-rata feature importance dari beberapa run.",
            "C) Cross-validation memberi estimasi performa lebih stabil dibanding 1 kali split.",
        ],
        font_size=17,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _set_slide_title(slide, "Modul 2: Contoh Kasus Nyata (Spam / Fraud)")
    _add_two_columns(
        slide,
        "Yang sering terjadi",
        [
            "Data tidak seimbang (fraud sangat sedikit)",
            "Accuracy bisa menipu (99% benar tapi semua diprediksi 'normal')",
            "Biaya salah beda: FN (miss fraud) biasanya lebih mahal",
        ],
        "Solusi evaluasi",
        [
            "Fokus metric: precision/recall/F1",
            "Gunakan confusion matrix untuk lihat FP vs FN",
            "Pertimbangkan threshold tuning",
            "Gunakan stratified split + cross-validation",
        ],
    )
    _add_note(slide, "Rule: pilih metric sesuai biaya bisnis, bukan sekadar accuracy.")

    # --- Modul 3 ---
    _add_section_slide(prs, "Modul 3", "Neural Network — MNIST")

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _set_slide_title(slide, "Modul 3: Tujuan & Output")
    _add_bullets_fs(
        slide,
        0.9,
        1.8,
        12.0,
        4.2,
        [
            "Tujuan: mengenali digit 0–9 dari gambar.",
            "Output: loss/accuracy per epoch + evaluasi test + prediksi contoh.",
            "File: 03_neural_network.py",
        ],
        font_size=18,
    )
    _add_code_box(slide, 0.9, 5.2, 12.0, 1.2, ['& "./.venv/Scripts/python.exe" 03_neural_network.py'])

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _set_slide_title(slide, "Modul 3: Komponen NN (Yang Harus Dihafal)")
    _add_two_columns(
        slide,
        "Komponen",
        [
            "Layer: Dense / Conv",
            "Activation: ReLU / Softmax",
            "Loss: cross-entropy",
            "Optimizer: Adam",
        ],
        "Makna",
        [
            "Layer = transformasi fitur",
            "Activation = non-linearitas",
            "Loss = ukuran salah",
            "Optimizer = cara update weights",
        ],
    )

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _set_slide_title(slide, "Modul 3: Mengapa Prediksi Bisa Benar?")
    _add_bullets_fs(
        slide,
        0.9,
        1.8,
        12.0,
        5.2,
        [
            "Output softmax = probabilitas 10 kelas.",
            "Backpropagation menghitung gradien: bagaimana weights harus berubah agar loss turun.",
            "Gradient descent/Adam melakukan update kecil berulang sampai model membaik.",
            "Overfitting muncul jika model terlalu kompleks / epoch terlalu banyak tanpa regularisasi.",
        ],
        font_size=18,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _set_slide_title(slide, "Modul 3: Latihan")
    _add_bullets_fs(
        slide,
        0.9,
        1.8,
        12.0,
        5.2,
        [
            "Latihan A: Ubah epochs (1, 3, 10). Catat train/test accuracy.",
            "Latihan B: Tambah Dropout. Apa dampaknya pada gap train vs val?",
            "Latihan C: Salah satu digit sering salah (mis. 4 vs 9). Kenapa? Cari contoh gambarnya.",
        ],
        font_size=18,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _set_slide_title(slide, "Modul 3: Pembahasan Latihan (Step-by-step)")
    _add_bullets_fs(
        slide,
        0.9,
        1.8,
        12.0,
        5.2,
        [
            "A) Epoch naik: train accuracy naik; kalau test/val tidak naik → mulai overfitting.",
            "B) Dropout: train accuracy bisa turun sedikit, tapi val/test bisa lebih stabil.",
            "C) Digit mirip (4 vs 9): bentuk tulisan bervariasi; pixel noise dan gaya tulisan bikin overlap.",
            "   Cara cek: tampilkan contoh yang salah → amati apakah memang ambigu.",
        ],
        font_size=17,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _set_slide_title(slide, "Modul 3: Contoh Kasus Nyata (OCR / Dokumen)")
    _add_two_columns(
        slide,
        "Masalah dunia nyata",
        [
            "Foto blur, miring, cahaya buruk",
            "Font beragam, tulisan tangan",
            "Background mengganggu",
        ],
        "Strategi",
        [
            "Preprocessing (crop, threshold, denoise)",
            "Data augmentation (rotasi, noise)",
            "Model CNN/CRNN untuk teks panjang",
            "Evaluasi: accuracy per karakter/word error rate",
        ],
    )
    _add_note(slide, "MNIST adalah latihan dasar. Kasus OCR nyata butuh data & preprocessing lebih serius.")

    # --- Modul 4 ---
    _add_section_slide(prs, "Modul 4", "Computer Vision — CNN CIFAR-10")

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _set_slide_title(slide, "Modul 4: Tujuan & Output")
    _add_bullets_fs(
        slide,
        0.9,
        1.8,
        12.0,
        4.2,
        [
            "Tujuan: klasifikasi gambar warna 10 kelas.",
            "Output: loss/accuracy + evaluasi + contoh prediksi.",
            "File: 04_computer_vision.py",
        ],
        font_size=18,
    )
    _add_code_box(slide, 0.9, 5.2, 12.0, 1.2, ['& "./.venv/Scripts/python.exe" 04_computer_vision.py'])

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _set_slide_title(slide, "Modul 4: Mengapa CNN Lebih Cocok untuk Gambar?")
    _add_bullets_fs(
        slide,
        0.9,
        1.8,
        12.0,
        5.2,
        [
            "Convolution menangkap pola lokal (edge/texture) tanpa harus 'melihat' seluruh gambar sekaligus.",
            "Pooling mereduksi ukuran dan membuat model lebih tahan terhadap pergeseran kecil.",
            "Layer awal belajar fitur sederhana; layer akhir belajar konsep objek.",
        ],
        font_size=18,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _set_slide_title(slide, "Modul 4: Latihan")
    _add_bullets_fs(
        slide,
        0.9,
        1.8,
        12.0,
        5.2,
        [
            "Latihan A: Tambah filter (mis. 16→32) dan lihat akurasi + waktu training.",
            "Latihan B: Tambah 1 layer conv. Apakah akurasi naik atau malah overfitting?",
            "Latihan C: Kurangi epoch untuk eksperimen cepat. Bagaimana trade-off waktu vs performa?",
        ],
        font_size=18,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _set_slide_title(slide, "Modul 4: Pembahasan Latihan (Step-by-step)")
    _add_bullets_fs(
        slide,
        0.9,
        1.8,
        12.0,
        5.2,
        [
            "A) Filter naik: kapasitas model naik → akurasi bisa naik, tapi training lebih lama.",
            "B) Layer bertambah: bisa menangkap pola lebih kompleks, tapi risiko overfitting naik.",
            "C) Epoch turun: cepat untuk eksperimen, tapi mungkin belum konvergen.",
            "   Praktik: cari baseline cepat dulu → baru training lebih lama saat arsitektur sudah oke.",
        ],
        font_size=17,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _set_slide_title(slide, "Modul 4: Contoh Kasus Nyata (Quality Check / Klasifikasi Produk)")
    _add_two_columns(
        slide,
        "Contoh",
        [
            "Deteksi cacat produk (retak/lecet)",
            "Klasifikasi jenis barang di gudang",
            "Klasifikasi makanan (fresh/tidak)",
        ],
        "Catatan penting",
        [
            "Data harus representatif (cahaya, kamera, sudut)",
            "Evaluasi: per kelas (kelas minoritas sering gagal)",
            "Waspada: domain shift (kamera produksi vs data training)",
            "Solusi: transfer learning + monitoring performa",
        ],
    )

    # --- Troubleshooting (lebih lengkap) ---
    _add_section_slide(prs, "Troubleshooting", "Jika ada error, cek ini dulu (urut)!")

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _set_slide_title(slide, "1) Pastikan Interpreter Benar")
    _add_bullets_fs(
        slide,
        0.9,
        1.8,
        12.0,
        3.0,
        [
            "Gejala: ModuleNotFoundError / versi package tidak sesuai.",
            "Solusi: jalankan dengan .venv dan cek sys.executable.",
        ],
        font_size=18,
    )
    _add_code_box(slide, 0.9, 4.1, 12.0, 1.3, ['& "./.venv/Scripts/python.exe" -c "import sys; print(sys.executable)"'])

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _set_slide_title(slide, "2) Install Dependency")
    _add_bullets_fs(
        slide,
        0.9,
        1.8,
        12.0,
        3.0,
        [
            "Jika pip install -r requirements.txt gagal: update pip dulu.",
            "Untuk kasus tertentu: install paket satu per satu untuk menemukan penyebab.",
        ],
        font_size=18,
    )
    _add_code_box(
        slide,
        0.9,
        4.1,
        12.0,
        1.6,
        [
            '& "./.venv/Scripts/python.exe" -m pip install --upgrade pip',
            '& "./.venv/Scripts/python.exe" -m pip install -r requirements.txt',
        ],
    )

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _set_slide_title(slide, "3) Program Terasa Lama")
    _add_bullets_fs(
        slide,
        0.9,
        1.8,
        12.0,
        5.2,
        [
            "Kemungkinan: download dataset (MNIST/CIFAR) pertama kali.",
            "Kemungkinan: training berjalan di CPU (wajar lebih lambat).",
            "Solusi cepat: kurangi epoch dulu untuk eksperimen.",
            "Solusi: pastikan storage cukup dan internet stabil.",
        ],
        font_size=18,
    )

    # --- Cara mencatat eksperimen (agar tidak bingung) ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _set_slide_title(slide, "Cara Mencatat Eksperimen (Wajib)")
    _add_bullets_fs(
        slide,
        0.9,
        1.8,
        12.0,
        5.2,
        [
            "Buat tabel sederhana: (Tanggal, Model, Setting, Metric, Catatan).",
            "Ubah 1 hal saja per eksperimen (mis. epoch, depth, filter).",
            "Simpan hasil: screenshot plot / log training / confusion matrix.",
            "Tujuan: Anda bisa menjawab 'perubahan apa yang membuat model membaik?'.",
        ],
        font_size=18,
    )

    # --- Evaluasi yang benar ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _set_slide_title(slide, "Checklist Evaluasi (Agar Tidak Salah Klaim)")
    _add_bullets_fs(
        slide,
        0.9,
        1.8,
        12.0,
        5.2,
        [
            "Apakah test set tidak pernah dipakai saat tuning?",
            "Apakah metric sesuai problem (dan biaya bisnis)?",
            "Apakah ada data leakage?",
            "Apakah Anda cek error cases (contoh yang salah)?",
            "Apakah performa stabil (cross-validation / beberapa seed)?",
        ],
        font_size=18,
    )

    # --- Deploy ringan ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _set_slide_title(slide, "Next Step: Deploy Sederhana (Opsional)")
    _add_two_columns(
        slide,
        "Pilihan mudah",
        [
            "Streamlit: UI cepat untuk demo",
            "FastAPI: API endpoint untuk model",
            "Notebook: laporan eksperimen",
        ],
        "Yang perlu dipikirkan",
        [
            "Simpan model (joblib/keras .h5/.keras)",
            "Input validation (hindari input aneh)",
            "Monitoring: data baru beda dari data training",
        ],
    )
    _add_note(slide, "Kalau Anda ingin, saya bisa bantu buat 1 demo Streamlit untuk Modul 1 atau Modul 2.")

    # --- Portofolio ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _set_slide_title(slide, "Portofolio (Agar Terlihat Progress)")
    _add_bullets_fs(
        slide,
        0.9,
        1.8,
        12.0,
        5.2,
        [
            "Buat 2–3 repo kecil: 1 regression, 1 classification, 1 vision.",
            "Setiap repo wajib ada: README, cara run, metric, contoh output.",
            "Tulis 'mengapa' Anda memilih metric/model tersebut.",
            "Ini yang membedakan 'sekadar run code' vs 'paham alur AI'.",
        ],
        font_size=18,
    )

    # --- Checklist project ---
    _add_section_slide(prs, "Next Level", "Kalau sudah selesai modul, buat 1 project nyata")

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _set_slide_title(slide, "Template Project (Wajib Ada)")
    _add_bullets_fs(
        slide,
        0.9,
        1.8,
        12.0,
        5.2,
        [
            "1) Problem statement (apa yang diprediksi, untuk siapa)",
            "2) Data: sumber, jumlah, fitur/target, pembersihan",
            "3) Baseline sederhana (model paling simple dulu)",
            "4) Evaluasi: metric yang tepat + interpretasi",
            "5) Improvement: 2–3 eksperimen dan catat hasil",
            "6) Kesimpulan: kapan model bisa dipakai, kapan tidak",
        ],
        font_size=18,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    _set_slide_title(slide, "Roadmap Latihan 2 Minggu")
    _add_flow(slide, ["Hari 1-2", "Hari 3-4", "Hari 5-7", "Minggu 2", "Project"])  # type: ignore[arg-type]
    _add_bullets_fs(
        slide,
        0.9,
        3.4,
        12.0,
        3.0,
        [
            "Hari 1–2: Regression + R² + visualisasi",
            "Hari 3–4: Classification + confusion matrix + depth tuning",
            "Hari 5–7: NN + epoch/batch + overfitting",
            "Minggu 2: CNN + eksperimen arsitektur sederhana",
            "Project: 1 kasus nyata + README hasil",
        ],
        font_size=18,
    )

    # --- Closing ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _set_slide_title(slide, "Penutup")
    _add_bullets_fs(
        slide,
        0.9,
        1.8,
        12.0,
        5.2,
        [
            "Skill inti AI developer pemula: bisa mengukur hasil (metric) dan menjelaskan 'mengapa'.",
            "Saran: ulangi tiap modul 2x sambil ubah parameter dan catat.",
            "Referensi: BUKU_PANDUAN_AI.pdf (buku) untuk penjelasan lebih panjang.",
        ],
        font_size=18,
    )

    # Final pass: apply comic background + font styling consistently
    for s in prs.slides:
        _apply_slide_background(s)
        if getattr(s.shapes, "title", None) is not None:
            _style_title_shape(s.shapes.title)
        for shape in s.shapes:
            if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
                _style_text_frame(shape.text_frame)

    return prs


def main() -> None:
    prs = build_deck()
    prs.save(OUT)
    print(f"✅ PPTX created: {OUT}")


if __name__ == "__main__":
    main()
